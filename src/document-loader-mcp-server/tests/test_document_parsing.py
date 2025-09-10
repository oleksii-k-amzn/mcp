# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
"""Test document parsing functionality by generating sample PDF, DOCX, XLSX, and PPTX files.

This module tests the MCP server tools against generated sample documents.
"""

import openpyxl

# Import required libraries for testing
import pdfplumber
import pytest
from docx import Document
from markitdown import MarkItDown
from openpyxl.styles import Font, PatternFill
from pathlib import Path
from pptx import Presentation

# Document generation imports
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer


def _read_pdf_helper(file_path: str) -> str:
    """Helper function to test PDF reading functionality."""
    try:
        text_content = ''
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text_content += f'\n--- Page {page_num} ---\n'
                page_text = page.extract_text()
                if page_text:
                    text_content += page_text
        return text_content.strip()
    except FileNotFoundError:
        return f'Error: Could not find PDF file at {file_path}'
    except Exception as e:
        return f'Error reading PDF file {file_path}: {str(e)}'


def _read_worddoc_helper(file_path: str) -> str:
    """Helper function to test Word document reading functionality."""
    try:
        md = MarkItDown()
        result = md.convert(file_path)
        return result.text_content
    except FileNotFoundError:
        return f'Error: Could not find Word document at {file_path}'
    except Exception as e:
        return f'Error reading Word document {file_path}: {str(e)}'


def _read_xlsx_helper(file_path: str) -> str:
    """Helper function to test Excel spreadsheet reading functionality."""
    try:
        md = MarkItDown()
        result = md.convert(file_path)
        return result.text_content
    except FileNotFoundError:
        return f'Error: Could not find Excel file at {file_path}'
    except Exception as e:
        return f'Error reading Excel file {file_path}: {str(e)}'


def _read_pptx_helper(file_path: str) -> str:
    """Helper function to test PowerPoint presentation reading functionality."""
    try:
        md = MarkItDown()
        result = md.convert(file_path)
        return result.text_content
    except FileNotFoundError:
        return f'Error: Could not find PowerPoint file at {file_path}'
    except Exception as e:
        return f'Error reading PowerPoint file {file_path}: {str(e)}'


@pytest.fixture(scope='session')
def document_generator():
    """Pytest fixture to provide a document generator instance."""
    return DocumentTestGenerator()


class DocumentTestGenerator:
    """Generate sample documents for testing."""

    def __init__(self, output_dir: str = 'tests/sample_docs'):
        """Initialize the sample document generator.

        Args:
            output_dir (str): Directory to store generated sample documents.
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def generate_sample_pdf(self) -> str:
        """Generate a sample PDF with various content types."""
        pdf_path = self.output_dir / 'sample_document.pdf'

        # Create PDF document
        doc = SimpleDocTemplate(str(pdf_path), pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        # Add title
        title = Paragraph('Sample PDF Document for Testing', styles['Title'])
        story.append(title)
        story.append(Spacer(1, 12))

        # Add heading
        heading = Paragraph('Introduction', styles['Heading1'])
        story.append(heading)
        story.append(Spacer(1, 6))

        # Add body text
        body_text = """
        This is a sample PDF document generated using ReportLab for testing the
        document parsing functionality of the Document Loader MCP server.

        The document contains various types of content including:
        • Headers and titles
        • Body paragraphs with formatting
        • Bullet points and lists
        • Multiple sections
        """

        body = Paragraph(body_text, styles['Normal'])
        story.append(body)
        story.append(Spacer(1, 12))

        # Add another section
        section_heading = Paragraph('Technical Details', styles['Heading1'])
        story.append(section_heading)
        story.append(Spacer(1, 6))

        technical_text = """
        This PDF parsing test validates that the pdfplumber library can successfully
        extract text content from generated PDF files. The test ensures that:

        1. Text extraction preserves content accuracy
        2. Formatting elements are handled appropriately
        3. Multi-paragraph documents are processed correctly
        4. Special characters and symbols are maintained
        """

        technical = Paragraph(technical_text, styles['Normal'])
        story.append(technical)

        # Build the PDF
        doc.build(story)
        return str(pdf_path)

    def generate_sample_docx(self) -> str:
        """Generate a sample DOCX with various content types."""
        docx_path = self.output_dir / 'sample_document.docx'

        # Create Word document
        doc = Document()

        # Add title
        doc.add_heading('Sample DOCX Document for Testing', 0)

        # Add introduction
        doc.add_heading('Introduction', level=1)
        doc.add_paragraph(
            'This is a sample DOCX document generated using python-docx for testing '
            'the document parsing functionality of the Document Loader MCP server.'
        )

        # Add bullet points
        doc.add_paragraph(
            'The document contains various types of content including:', style='Normal'
        )
        doc.add_paragraph('Headers and titles', style='List Bullet')
        doc.add_paragraph('Body paragraphs with formatting', style='List Bullet')
        doc.add_paragraph('Bullet points and lists', style='List Bullet')
        doc.add_paragraph('Multiple sections', style='List Bullet')

        # Add another section
        doc.add_heading('Technical Implementation', level=1)
        doc.add_paragraph(
            'This DOCX parsing test validates that the markitdown library can successfully '
            'extract and convert content from generated Word documents. The test ensures that:'
        )

        # Add numbered list
        doc.add_paragraph('Text extraction preserves content accuracy', style='List Number')
        doc.add_paragraph(
            'Formatting elements are converted to markdown appropriately', style='List Number'
        )
        doc.add_paragraph('Multi-paragraph documents are processed correctly', style='List Number')
        doc.add_paragraph(
            'Headers and structure are maintained in markdown format', style='List Number'
        )

        # Add conclusion
        doc.add_heading('Conclusion', level=1)
        doc.add_paragraph(
            'This test document provides a comprehensive example for validating '
            'document parsing capabilities across different content types and structures.'
        )

        # Save the document
        doc.save(str(docx_path))
        return str(docx_path)

    def generate_sample_xlsx(self) -> str:
        """Generate a sample Excel file with multiple sheets and data."""
        xlsx_path = self.output_dir / 'sample_spreadsheet.xlsx'

        # Create Excel workbook
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = 'Test Data'

        # Add headers
        headers = ['Name', 'Age', 'Department', 'Salary', 'Start Date']
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            if cell:
                cell.font = Font(bold=True)
                cell.fill = PatternFill(
                    start_color='CCCCCC', end_color='CCCCCC', fill_type='solid'
                )

        # Add test data
        test_data = [
            ['Alice Johnson', 28, 'Engineering', 75000, '2022-01-15'],
            ['Bob Smith', 32, 'Marketing', 65000, '2021-03-10'],
            ['Carol Davis', 45, 'Sales', 80000, '2020-07-22'],
            ['David Brown', 29, 'HR', 55000, '2023-02-01'],
            ['Eve Wilson', 38, 'Finance', 70000, '2019-11-05'],
        ]

        for row, record in enumerate(test_data, 2):
            for col, value in enumerate(record, 1):
                cell = ws.cell(row=row, column=col, value=value)

        # Add a summary sheet
        ws2 = wb.create_sheet('Summary')
        ws2['A1'] = 'Department Summary Report'
        ws2['A1'].font = Font(bold=True, size=14)

        ws2['A3'] = 'Department'
        ws2['B3'] = 'Employee Count'
        ws2['C3'] = 'Average Salary'

        # Make headers bold
        for cell in ['A3', 'B3', 'C3']:
            ws2[cell].font = Font(bold=True)

        summary_data = [
            ['Engineering', 1, 75000],
            ['Marketing', 1, 65000],
            ['Sales', 1, 80000],
            ['HR', 1, 55000],
            ['Finance', 1, 70000],
        ]

        for row, record in enumerate(summary_data, 4):
            for col, value in enumerate(record, 1):
                ws2.cell(row=row, column=col, value=value)

        # Save the workbook
        wb.save(str(xlsx_path))
        return str(xlsx_path)

    def generate_sample_pptx(self) -> str:
        """Generate a sample PowerPoint presentation."""
        pptx_path = self.output_dir / 'sample_presentation.pptx'

        # Create PowerPoint presentation
        prs = Presentation()

        # Slide 1: Title slide
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        if title and hasattr(title, 'text'):
            title.text = 'Test Presentation'
        if subtitle and hasattr(subtitle, 'text'):
            subtitle.text = 'Document Loader MCP Server Testing\nGenerated for validation purposes'

        # Slide 2: Content slide
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        if title and hasattr(title, 'text'):
            title.text = 'Testing Features'
        if content and hasattr(content, 'text'):
            content.text = """• PDF document parsing
• Word document processing
• Excel spreadsheet conversion
• PowerPoint presentation extraction
• Markdown output generation
• Comprehensive error handling"""

        # Slide 3: Data validation slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        if title and hasattr(title, 'text'):
            title.text = 'Test Data Validation'
        if content and hasattr(content, 'text'):
            content.text = """Validation Criteria:
• Content extraction accuracy
• Format preservation
• Multi-sheet/slide support
• Error handling robustness
• Markdown conversion quality"""

        # Slide 4: Results
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        if title and hasattr(title, 'text'):
            title.text = 'Expected Results'
        if content and hasattr(content, 'text'):
            content.text = """This test validates:
✓ Multi-format document support
✓ Structured content extraction
✓ Reliable markdown conversion
✓ MCP protocol integration

Test completion indicates successful functionality!"""

        # Save the presentation
        prs.save(str(pptx_path))
        return str(pptx_path)

    def generate_sample_image(self) -> str:
        """Generate a sample image file for testing."""
        from PIL import Image, ImageDraw, ImageFont

        # Create a simple test image
        image_path = self.output_dir / 'sample_image.png'

        # Create a 400x300 image with a light blue background
        img = Image.new('RGB', (400, 300), color='lightblue')
        draw = ImageDraw.Draw(img)

        # Add some text and shapes
        try:
            # Try to use a default font, fall back to basic if not available
            font = ImageFont.load_default()
        except Exception:
            font = None

        # Draw some shapes and text
        draw.rectangle([50, 50, 350, 100], fill='white', outline='black', width=2)
        draw.text((60, 65), 'Document Loader MCP Server', fill='black', font=font)
        draw.text((60, 85), 'Test Image Generation', fill='black', font=font)

        # Draw some geometric shapes
        draw.ellipse([50, 120, 150, 220], fill='yellow', outline='orange', width=3)
        draw.rectangle([200, 120, 300, 220], fill='lightgreen', outline='darkgreen', width=3)
        draw.polygon([(325, 120), (375, 170), (325, 220), (275, 170)], fill='pink', outline='red')

        # Add some test information
        draw.text((50, 240), 'Generated for MCP testing', fill='darkblue', font=font)
        draw.text(
            (50, 260), f'Size: {img.size[0]}x{img.size[1]} pixels', fill='darkblue', font=font
        )

        # Save the image
        img.save(str(image_path), 'PNG')
        return str(image_path)


def test_pdf_parsing(document_generator):
    """Test PDF parsing functionality."""
    print('Testing PDF parsing...')

    pdf_path = document_generator.generate_sample_pdf()
    print(f'Generated PDF: {pdf_path}')

    # Ensure the PDF file was created
    assert Path(pdf_path).exists(), f'PDF file should exist at {pdf_path}'

    # Test the read_pdf tool
    text_content = _read_pdf_helper(pdf_path)

    # Verify we got content back
    assert text_content, 'PDF parsing should return non-empty content'
    assert not text_content.startswith('Error:'), f'PDF parsing failed: {text_content}'

    print('PDF parsing successful!')
    print(f'Extracted text length: {len(text_content)} characters')
    print(f'Content preview: {text_content[:100]}...')

    # Verify key content is present
    assert 'Sample PDF Document for Testing' in text_content
    assert 'Introduction' in text_content
    assert 'Technical Details' in text_content
    assert 'pdfplumber library' in text_content

    print('✓ PDF content validation passed')


def test_docx_parsing(document_generator):
    """Test DOCX parsing functionality."""
    print('Testing DOCX parsing...')

    docx_path = document_generator.generate_sample_docx()
    print(f'Generated DOCX: {docx_path}')

    # Ensure the DOCX file was created
    assert Path(docx_path).exists(), f'DOCX file should exist at {docx_path}'

    # Test the read_worddoc tool
    markdown_content = _read_worddoc_helper(docx_path)

    # Verify we got content back
    assert markdown_content, 'DOCX parsing should return non-empty content'
    assert not markdown_content.startswith('Error:'), f'DOCX parsing failed: {markdown_content}'

    print('DOCX parsing successful!')
    print(f'Extracted markdown length: {len(markdown_content)} characters')

    # Show a preview of the extracted content
    print(f'Content preview: {markdown_content[:100]}...')

    # Verify key content is present
    assert 'Sample DOCX Document for Testing' in markdown_content
    assert 'Introduction' in markdown_content
    assert 'Technical Implementation' in markdown_content
    assert 'markitdown library' in markdown_content

    print('✓ DOCX content validation passed')


def test_xlsx_parsing(document_generator):
    """Test Excel spreadsheet parsing functionality."""
    print('Testing XLSX parsing...')

    xlsx_path = document_generator.generate_sample_xlsx()
    print(f'Generated XLSX: {xlsx_path}')

    # Ensure the XLSX file was created
    assert Path(xlsx_path).exists(), f'XLSX file should exist at {xlsx_path}'

    # Test the read_xlsx tool
    markdown_content = _read_xlsx_helper(xlsx_path)

    # Verify we got content back
    assert markdown_content, 'XLSX parsing should return non-empty content'
    assert not markdown_content.startswith('Error:'), f'XLSX parsing failed: {markdown_content}'

    print('XLSX parsing successful!')
    print(f'Extracted markdown length: {len(markdown_content)} characters')

    # Show a preview of the extracted content
    print(f'Content preview: {markdown_content[:200]}...')

    # Verify key content is present
    assert 'Test Data' in markdown_content or 'Name' in markdown_content
    assert 'Alice Johnson' in markdown_content
    assert 'Engineering' in markdown_content
    assert 'Department Summary' in markdown_content

    print('✓ XLSX content validation passed')


def test_pptx_parsing(document_generator):
    """Test PowerPoint presentation parsing functionality."""
    print('Testing PPTX parsing...')

    pptx_path = document_generator.generate_sample_pptx()
    print(f'Generated PPTX: {pptx_path}')

    # Ensure the PPTX file was created
    assert Path(pptx_path).exists(), f'PPTX file should exist at {pptx_path}'

    # Test the read_pptx tool
    markdown_content = _read_pptx_helper(pptx_path)

    # Verify we got content back
    assert markdown_content, 'PPTX parsing should return non-empty content'
    assert not markdown_content.startswith('Error:'), f'PPTX parsing failed: {markdown_content}'

    print('PPTX parsing successful!')
    print(f'Extracted markdown length: {len(markdown_content)} characters')

    # Show a preview of the extracted content
    print(f'Content preview: {markdown_content[:200]}...')

    # Verify key content is present
    assert 'Test Presentation' in markdown_content
    assert 'Testing Features' in markdown_content
    assert 'PDF document parsing' in markdown_content
    assert 'Expected Results' in markdown_content

    print('✓ PPTX content validation passed')


def test_sample_documents_exist(document_generator):
    """Test that sample documents are generated and exist."""
    pdf_path = document_generator.generate_sample_pdf()
    docx_path = document_generator.generate_sample_docx()
    xlsx_path = document_generator.generate_sample_xlsx()
    pptx_path = document_generator.generate_sample_pptx()

    assert Path(pdf_path).exists(), f'PDF file should exist at {pdf_path}'
    assert Path(docx_path).exists(), f'DOCX file should exist at {docx_path}'
    assert Path(xlsx_path).exists(), f'XLSX file should exist at {xlsx_path}'
    assert Path(pptx_path).exists(), f'PPTX file should exist at {pptx_path}'

    # Check file sizes are reasonable
    pdf_size = Path(pdf_path).stat().st_size
    docx_size = Path(docx_path).stat().st_size
    xlsx_size = Path(xlsx_path).stat().st_size
    pptx_size = Path(pptx_path).stat().st_size

    assert pdf_size > 1000, f'PDF file seems too small: {pdf_size} bytes'
    assert docx_size > 1000, f'DOCX file seems too small: {docx_size} bytes'
    assert xlsx_size > 1000, f'XLSX file seems too small: {xlsx_size} bytes'
    assert pptx_size > 1000, f'PPTX file seems too small: {pptx_size} bytes'

    print('✓ Sample documents exist and have reasonable sizes')
    print(f'  PDF: {pdf_size} bytes')
    print(f'  DOCX: {docx_size} bytes')
    print(f'  XLSX: {xlsx_size} bytes')
    print(f'  PPTX: {pptx_size} bytes')
